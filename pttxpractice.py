import copy
import pptx
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.fill import FillFormat
from lxml import etree
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

def duplicateSlide(new_slide, old_slide):
    sp = new_slide.shapes.title.element
    sp.getparent().remove(sp)
    sp = new_slide.shapes.placeholders[1].element
    sp.getparent().remove(sp)
    for shp in old_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return new_slide

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_cell_border(cell, border_color='000000', border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

if __name__ == "__main__":
    pptx_path = "D:/[7]研發簡報/週會/Test_Report.pptx"
    image_path = "D:/[3]DeepLearning/[2]Model/[2]Keras/[2]CNN/model_V2.10.1/model_b2_e100_d03_1/"
    line_pt_space = Pt(10)
    add_weight = False
    draw_table = True
    draw_plot = True
    loss_weight_p = str(0.2)
    loss_weight_n = str(0.8)
    prs = Presentation(pptx_path)
    bullet_slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

#Title
    title_shape = shapes.title
    title_shape.left = Cm(1.86)
    title_shape.top = Cm(0)
    title_shape.width = Cm(21.68)
    title_shape.height = Cm(3.23)
    tf = title_shape.text_frame

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Deep Learning"
    p.font.size = Pt(44)
    
#Content
    body_shape = shapes.placeholders[1]
    body_shape.left = Cm(1.32)
    body_shape.top = Cm(1.48)
    body_shape.width = Cm(23)
    body_shape.height = Cm(15.85)
    tf = body_shape.text_frame

    p = tf.paragraphs[0]
    p.line_spacing = line_pt_space
    p.level = 0
    p.text = 'CNN'
    p.font.size = Pt(28)

    p = tf.add_paragraph()
    p.line_spacing = line_pt_space
    p.level = 1
    p.text = "Model Test"
    p.font.size = Pt(24)

    p = tf.add_paragraph()
    p.line_spacing = line_pt_space
    p.level = 2
    p.text = "Input: 1644 images (All Transverse)"
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.line_spacing = line_pt_space
    p.level = 2
    p.text = "Node: 256, Batch Size: 2, Epochs: 100, Drop Out: 0.8"
    p.font.size = Pt(18)

#Region: Convert MathML (MML) into Office MathML (OMML) using a XSLT stylesheet
    if add_weight:
        mathml = '<math xmlns="http://www.w3.org/1998/Math/MathML"><mi>l</mi><mi>o</mi><mi>s</mi><mi>s</mi><mo>&#xA0;</mo><mo>=</mo><mo>&#xA0;</mo><mo>-</mo><mfenced><mrow><mn>'+loss_weight_p+'</mn><mo>&#xD7;</mo><mi>y</mi><mo>&#xD7;</mo><msub><mi>log</mi><mi>e</mi></msub><mi>p</mi><mo>+</mo><mn>'+loss_weight_n+'</mn><mo>&#xD7;</mo><mfenced><mrow><mn>1</mn><mo>-</mo><mi>y</mi></mrow></mfenced><mo>&#xD7;</mo><msub><mi>log</mi><mi>e</mi></msub><mfenced><mrow><mn>1</mn><mo>-</mo><mi>p</mi></mrow></mfenced></mrow></mfenced></math>'
    else:
        mathml = '<math xmlns="http://www.w3.org/1998/Math/MathML"><mi>l</mi><mi>o</mi><mi>s</mi><mi>s</mi><mo>&#xA0;</mo><mo>=</mo><mo>&#xA0;</mo><mo>-</mo><mfenced><mrow><mi>y</mi><mo>&#xD7;</mo><msub><mi>log</mi><mi>e</mi></msub><mi>p</mi><mo>+</mo><mfenced><mrow><mn>1</mn><mo>-</mo><mi>y</mi></mrow></mfenced><mo>&#xD7;</mo><msub><mi>log</mi><mi>e</mi></msub><mfenced><mrow><mn>1</mn><mo>-</mo><mi>p</mi></mrow></mfenced></mrow></mfenced></math>'   
    tree = etree.fromstring(mathml)
    xslt = etree.parse('C:/Program Files/Microsoft Office 15/root/office15/MML2OMML.XSL')

    wrapper = etree.fromstring('<a14:m xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main"><m:oMathPara xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math"></m:oMathPara></a14:m>')

    transform = etree.XSLT(xslt)
    new_dom = transform(tree)

    wrapper.getchildren()[0].append(new_dom.getroot())

    p = tf.add_paragraph()
    p.line_spacing = line_pt_space
    p.level = 2
    p._element.append(wrapper)
    p.font.size = Pt(18)
#End Region

    p = tf.add_paragraph()
    p.line_spacing = line_pt_space
    p.level = 1
    p.text = "Prediction Result"
    p.font.size = Pt(24)

#Table
    if draw_table:
        x, y, cx, cy = Cm(2.65), Cm(6.99), Cm(20.11), Cm(3.6)
        shape = slide.shapes.add_table(4, 5, x, y, cx, cy)
        #shape.table.fill.background()
        table = shape.table

        input_matrix = ([["model_b2_e100_d05_1","Loss", "Accuracy", "Sensitivity", "Specificity"],
                        ["Training", "0", "0", "0", "0"],
                        ["Validation", "0", "0", "0", "0"],
                        ["Testing", "0", "0", "0", "0"]])

        size_matrix = ([[12,12, 12, 12, 12],
                        [12, 16, 16, 16, 16],
                        [12, 16, 16, 16, 16],
                        [12, 16, 16, 16, 16]])

        table.columns[0].width = Cm(5.55)

        for (i,j) in [(i,j) for i in range(4) for j in range(5)]:
            cell = table.cell(i,j)
            cell.text = input_matrix[i][j]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(size_matrix[i][j])
            p.font.color.rgb = RGBColor(0,0,0)
            _set_cell_border(cell) 
            table.cell(i,j).fill.background()
    
    if draw_plot:
        second_slide = prs.slides.add_slide(slide.slide_layout)
        second_slide = duplicateSlide(second_slide, slide)
        current_shapes=second_slide.shapes
        tf = current_shapes.placeholders[1].text_frame
        p = tf.paragraphs[5]
        p.clear()
        sp = current_shapes[2].element
        sp.getparent().remove(sp)

        picture = current_shapes.add_picture(image_path+'model_b2_e100_d03_1_loss.png', Cm(4.33), Cm(6.35), Cm(16), Cm(12))

        pass
    prs.save(pptx_path)
    print('Done')